from __future__ import annotations

import csv
import json
import mimetypes
from pathlib import Path
from typing import Any

MAX_TEXT_CHARS = 500_000
MAX_PREVIEW_ROWS = 30
MAX_PREVIEW_COLS = 30


def _truncate(text: str) -> str:
    return text[:MAX_TEXT_CHARS]


def _safe_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk"):
        try:
            return _truncate(path.read_text(encoding=encoding))
        except UnicodeDecodeError:
            continue
    return _truncate(path.read_text(encoding="utf-8", errors="ignore"))


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    chunks: list[str] = []
    for page in reader.pages[:300]:
        chunks.append(page.extract_text() or "")
        if sum(len(x) for x in chunks) >= MAX_TEXT_CHARS:
            break
    return _truncate("\n".join(chunks))


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    chunks = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            chunks.append("\t".join(cell.text for cell in row.cells))
    return _truncate("\n".join(chunks))


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = getattr(shape, "text", "")
                if text:
                    chunks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    chunks.append("\t".join(cell.text for cell in row.cells))
    return _truncate("\n".join(chunks))


def _extract_xlsx(path: Path) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    text_chunks: list[str] = []
    preview: list[list[Any]] = []
    schema: list[str] = []
    total_rows = 0
    max_cols = 0

    try:
        for ws_index, ws in enumerate(wb.worksheets):
            text_chunks.append(f"[Sheet: {ws.title}]")
            sheet_rows = 0
            for row in ws.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row[:MAX_PREVIEW_COLS]]
                if not any(values):
                    continue
                total_rows += 1
                sheet_rows += 1
                max_cols = max(max_cols, len(values))
                text_chunks.append("\t".join(values))
                if ws_index == 0 and len(preview) < MAX_PREVIEW_ROWS:
                    preview.append(values)
                if sum(len(x) for x in text_chunks) >= MAX_TEXT_CHARS:
                    break
            if ws_index == 0 and preview:
                schema = [str(x) or f"column_{i+1}" for i, x in enumerate(preview[0])]
            if sum(len(x) for x in text_chunks) >= MAX_TEXT_CHARS:
                break
    finally:
        wb.close()

    return _truncate("\n".join(text_chunks)), {
        "rows": max(total_rows - 1, 0),
        "columns": max_cols,
        "schema": schema,
        "preview": preview,
    }


def _extract_csv(path: Path) -> tuple[str, dict[str, Any]]:
    sample = path.read_bytes()[:100_000]
    encoding = "utf-8-sig"
    try:
        sample.decode(encoding)
    except UnicodeDecodeError:
        encoding = "gb18030"

    text_chunks: list[str] = []
    preview: list[list[str]] = []
    total_rows = 0
    max_cols = 0
    schema: list[str] = []

    with path.open("r", encoding=encoding, errors="ignore", newline="") as fh:
        try:
            dialect = csv.Sniffer().sniff(fh.read(8192))
            fh.seek(0)
        except csv.Error:
            fh.seek(0)
            dialect = csv.excel
        reader = csv.reader(fh, dialect)
        for row in reader:
            values = [str(v) for v in row[:MAX_PREVIEW_COLS]]
            total_rows += 1
            max_cols = max(max_cols, len(values))
            text_chunks.append("\t".join(values))
            if len(preview) < MAX_PREVIEW_ROWS:
                preview.append(values)
            if len("\n".join(text_chunks)) >= MAX_TEXT_CHARS:
                break

    if preview:
        schema = [x or f"column_{i+1}" for i, x in enumerate(preview[0])]
    return _truncate("\n".join(text_chunks)), {
        "rows": max(total_rows - 1, 0),
        "columns": max_cols,
        "schema": schema,
        "preview": preview,
    }


def _extract_json(path: Path) -> tuple[str, dict[str, Any]]:
    raw = _safe_text(path)
    metadata = {"rows": None, "columns": None, "schema": [], "preview": []}
    try:
        data = json.loads(raw)
        if isinstance(data, list):
            metadata["rows"] = len(data)
            if data and isinstance(data[0], dict):
                metadata["schema"] = list(data[0].keys())[:MAX_PREVIEW_COLS]
                metadata["columns"] = len(metadata["schema"])
                metadata["preview"] = [
                    [item.get(k, "") for k in metadata["schema"]]
                    for item in data[:MAX_PREVIEW_ROWS]
                    if isinstance(item, dict)
                ]
        elif isinstance(data, dict):
            metadata["schema"] = list(data.keys())[:MAX_PREVIEW_COLS]
            metadata["columns"] = len(metadata["schema"])
            metadata["rows"] = 1
            metadata["preview"] = [[data.get(k, "") for k in metadata["schema"]]]
    except json.JSONDecodeError:
        pass
    return raw, metadata


def _image_metadata(path: Path) -> dict[str, Any]:
    from PIL import Image

    with Image.open(path) as image:
        return {
            "rows": int(image.height),
            "columns": int(image.width),
            "schema": ["width", "height", "format", "mode"],
            "preview": [[image.width, image.height, image.format or "", image.mode]],
        }


def extract_file(path: Path) -> dict[str, Any]:
    suffix = path.suffix.lower()
    mime_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    result: dict[str, Any] = {
        "content": "",
        "mime_type": mime_type,
        "rows": None,
        "columns": None,
        "schema": [],
        "preview": [],
        "error": "",
    }

    try:
        if suffix in {".txt", ".md", ".py", ".r", ".m", ".log", ".yaml", ".yml", ".ini", ".cfg", ".tex"}:
            result["content"] = _safe_text(path)
        elif suffix == ".pdf":
            result["content"] = _extract_pdf(path)
        elif suffix == ".docx":
            result["content"] = _extract_docx(path)
        elif suffix == ".pptx":
            result["content"] = _extract_pptx(path)
        elif suffix in {".xlsx", ".xlsm"}:
            result["content"], meta = _extract_xlsx(path)
            result.update(meta)
        elif suffix == ".csv":
            result["content"], meta = _extract_csv(path)
            result.update(meta)
        elif suffix == ".json":
            result["content"], meta = _extract_json(path)
            result.update(meta)
        elif suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp"}:
            result.update(_image_metadata(path))
        else:
            result["content"] = ""
    except Exception as exc:  # Upload should survive extraction failures.
        result["error"] = f"文本提取失败：{exc}"

    return result
